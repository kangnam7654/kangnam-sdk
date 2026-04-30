#!/usr/bin/env python3
"""Read a .pptx with python-pptx and print structural stats as JSON.

Invoked by the Rust integration test. Output example:

    {"slides":1,"text_boxes":1,"shapes":0,"images":0,"first_text":"hi"}

Exit code 0 = success; non-zero = parse failure.
"""
import json, sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def main(path: str) -> int:
    pres = Presentation(path)
    total_text = 0
    total_shapes = 0
    total_images = 0
    first_text = ""
    for slide in pres.slides:
        for sh in slide.shapes:
            if sh.has_text_frame and sh.text_frame.text:
                total_text += 1
                if not first_text:
                    first_text = sh.text_frame.text
            elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
            else:
                total_shapes += 1
    print(json.dumps({
        "slides": len(pres.slides),
        "text_boxes": total_text,
        "shapes": total_shapes,
        "images": total_images,
        "first_text": first_text,
    }))
    return 0

if __name__ == "__main__":
    sys.exit(main(sys.argv[1]))
